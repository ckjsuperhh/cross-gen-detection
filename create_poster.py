"""
Generate academic poster based on CS3308-ML-poster-template.
Fixed version: smaller fonts, concise text, proper layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Load template
prs = Presentation("CS3308-ML-poster-template(1).pptx")
slide = prs.slides[0]

def update_text(shape, new_text, font_size=None, bold=None):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = new_text
    if font_size and p.runs:
        p.runs[0].font.size = Pt(font_size)
    if bold is not None and p.runs:
        p.runs[0].font.bold = bold

def update_multiline(shape, lines, font_size=16, bold=False):
    """Add multiple lines to a text box."""
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            p.runs[0].font.size = Pt(font_size)
            if bold:
                p.runs[0].font.bold = True

# Shape mapping
shape_map = {}
for shape in slide.shapes:
    name = shape.name
    if "TextBox 26" in name:
        shape_map["title"] = shape
    elif "Rectangle 6" in name and shape.text and "Yuchen Wang" in shape.text:
        shape_map["author"] = shape
    elif "TextBox 2" in name and "Real-time" in shape.text:
        shape_map["section1"] = shape
    elif "TextBox 46" in name:
        shape_map["section2"] = shape
    elif "TextBox 7" in name:
        shape_map["intro"] = shape
    elif "Rectangle 12" in name and "Terminologies" in shape.text:
        shape_map["terminology_title"] = shape
    elif "TextBox 60" in name:
        shape_map["task"] = shape
    elif "Rectangle 61" in name:
        shape_map["challenges_title"] = shape
    elif "TextBox 62" in name:
        shape_map["challenges"] = shape
    elif "TextBox 93" in name:
        shape_map["method_title"] = shape
    elif "TextBox 95" in name:
        shape_map["method_desc"] = shape
    elif "TextBox 123" in name:
        shape_map["method2_title"] = shape
    elif "Rectangle 125" in name:
        shape_map["experiments_title"] = shape
    elif "TextBox 128" in name:
        shape_map["experiments_text"] = shape

print("Updating poster content (fixed version)...\n")

# Title (24.60in × 2.29in - smaller font)
if "title" in shape_map:
    update_text(shape_map["title"],
                "Why Detectors Fail on Unseen Generators",
                font_size=36, bold=True)
    print("  ✓ Title")

# Author
if "author" in shape_map:
    update_text(shape_map["author"],
                "Anonymous Author | AI-Generated Image Detection Project",
                font_size=18, bold=False)
    print("  ✓ Author")

# Section 1: Introduction
if "section1" in shape_map:
    update_text(shape_map["section1"], "Introduction", font_size=24, bold=True)
    print("  ✓ Section 1 title")

# Intro text (14.69in × 4.31in)
if "intro" in shape_map:
    lines = [
        "AI-generated images from GANs, diffusion models, and autoregressive",
        "architectures are increasingly photorealistic, creating urgent demand",
        "for robust detectors.",
        "",
        "Core challenge: cross-generator generalisation",
        "A detector trained on one generator typically fails on a different",
        "generator because each model leaves distinct low-level artifacts.",
        "",
        "Two strategies:",
        "(1) Data diversity: train on many generators",
        "(2) Representation universality: use frozen pretrained features",
        "",
        "No single method generalises robustly across all generator families.",
    ]
    update_multiline(shape_map["intro"], lines, font_size=16)
    print("  ✓ Intro text")

# Section 2: Results
if "section2" in shape_map:
    update_text(shape_map["section2"], "Key Findings", font_size=24, bold=True)
    print("  ✓ Section 2 title")

# Terminology → Dataset
if "terminology_title" in shape_map:
    update_text(shape_map["terminology_title"], "Dataset", font_size=22, bold=True)
    print("  ✓ Terminology title")

# Task (17.20in × 0.77in - very short height, must be concise)
if "task" in shape_map:
    lines = [
        "GenImage subset: GLIDE (train), BigGAN/ADM/VQDM (test)",
        "Backbones: ResNet-18, ViT-Small, CLIP-ViT-B/32",
        "Strategies: Full FT, Linear Probe, KNN",
    ]
    update_multiline(shape_map["task"], lines, font_size=14)
    print("  ✓ Task description")

# Challenges → Findings
if "challenges_title" in shape_map:
    update_text(shape_map["challenges_title"], "5 Key Findings", font_size=22, bold=True)
    print("  ✓ Challenges title")

# Challenges text (15.52in × 2.91in)
if "challenges" in shape_map:
    lines = [
        "1. Full FT dominates: 98–99% on GLIDE",
        "2. Diffusion collapse: ADM/VQDM 52–65%",
        "3. CLIP KNN best: 64.9% ADM, 64.7% VQDM",
        "4. Patch errors systematic: 90% consistency",
        "5. Features encode generators: Purity=0.559",
    ]
    update_multiline(shape_map["challenges"], lines, font_size=16)
    print("  ✓ Challenges text")

# Method title → Main Results
if "method_title" in shape_map:
    update_text(shape_map["method_title"],
                "Main Results: Detection Accuracy (%)",
                font_size=24, bold=True)
    print("  ✓ Method title")

# Method desc (8.61in × 6.41in)
if "method_desc" in shape_map:
    lines = [
        "Method       Backbone    Strategy      GLIDE  BigGAN  ADM   VQDM",
        "─────────────────────────────────────────────────────────────────",
        "Baseline     ResNet-18   Full FT       98.2   99.3    56.2  53.0",
        "Ablation     ViT-Small   Full FT       98.7   97.8    54.9  52.1",
        "Ablation     ResNet-18   Linear Probe  92.0   91.9    59.4  54.7",
        "Ablation     CLIP-ViT    Linear Probe  99.3   98.3    61.3  55.6",
        "Ojha-style   CLIP-ViT    KNN (k=1)     85.4   81.7    64.9  64.7",
        "Patch        ResNet-18   3×3 Vote      61.5   54.0    60.8  61.8",
        "",
        "GLIDE: in-distribution; others: zero-shot cross-generator",
    ]
    update_multiline(shape_map["method_desc"], lines, font_size=14)
    print("  ✓ Method description")

# Method 2 title → MoE
if "method2_title" in shape_map:
    update_text(shape_map["method2_title"],
                "MoE Framework (Future Work)",
                font_size=24, bold=True)
    print("  ✓ Method 2 title")

# Experiments → Clustering
if "experiments_title" in shape_map:
    update_text(shape_map["experiments_title"], "Clustering Evidence", font_size=22, bold=True)
    print("  ✓ Experiments title")

# Experiments text (15.23in × 9.22in)
if "experiments_text" in shape_map:
    lines = [
        "Unsupervised K-Means (k=4) on 800 fake images:",
        "",
        "Feature Space            Purity   NMI",
        "──────────────────────────────────────",
        "Pretrained ViT-Small     0.278    0.003",
        "Fine-tuned ViT-Small     0.559    0.378",
        "",
        "Cluster breakdown (fine-tuned):",
        "  Cluster 0: GLIDE (192/318)",
        "  Cluster 1+3: ADM/VQDM (154/312)",
        "  Cluster 2: BigGAN (97/178)",
        "",
        "Key insight: detection-oriented supervision",
        "implicitly structures feature space by",
        "generator family, enabling MoE routing.",
        "",
        "Proposed MoE architecture:",
        "(1) GAN expert — CNN for frequency artifacts",
        "(2) Diffusion expert — DIRE/FIRE detector",
        "(3) Universal router — routes by generator family",
    ]
    update_multiline(shape_map["experiments_text"], lines, font_size=16)
    print("  ✓ Experiments text")

# Save
output_path = "/root/ml-project/poster_ai_detection.pptx"
prs.save(output_path)
print(f"\n✓ Poster saved to {output_path}")
print(f"  Size: {prs.slide_width.inches:.1f} × {prs.slide_height.inches:.1f} inches")
print("\nTips:")
print("- Open in PowerPoint for best viewing")
print("- Adjust font sizes if needed")
print("- Add figures from outputs_ablation/figures/")
